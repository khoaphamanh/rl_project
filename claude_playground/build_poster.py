"""Build poster_tbptt.pptx: the LUH A0 template filled with this project.

Run from this directory:  python build_poster.py

Sections are auto-sized: content is laid out first with measured text heights,
then the panel frame is drawn around it and the columns are distributed over
the page.
"""

import os
import re
from PIL import ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "poster.pptx")
OUT = os.path.join(HERE, "poster_tbptt.pptx")
FIG = os.path.join(HERE, "figures") + os.sep
# compare_curves.png is the rasterised compare_curves.pdf (PowerPoint
# cannot place a PDF); qr_repo.png points at the repository.
# model_architecture_poster.png is the rasterised .svg of the same name, and
# ppo_losses.png is written by _make_loss_figure() below.

BLUE = RGBColor(0x0B, 0x53, 0x94)
DARK = RGBColor(0x20, 0x21, 0x24)
MUTED = RGBColor(0x5F, 0x63, 0x68)
BORDER = RGBColor(0x8A, 0x8A, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x13, 0x73, 0x33)
RED = RGBColor(0xC5, 0x22, 0x1F)
ORANGE = RGBColor(0xB0, 0x5A, 0x00)
TINT = RGBColor(0xEE, 0xF3, 0xFA)
GREY = RGBColor(0xF3, 0xF3, 0xF3)
SILVER = RGBColor(0x9E, 0x9E, 0x9E)

SANS = "IBM Plex Sans"
MONO = "IBM Plex Mono"
# any Arial-metric font will do: it is only used to measure line breaks
ARIAL = next((p for p in (
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
) if os.path.exists(p)))
LH = 1.30           # line height as a multiple of the point size
SAFETY = 1.07       # IBM Plex Sans runs a little wider than the Arial we measure with

prs = Presentation(SRC)
slide = prs.slides[0]
shapes = slide.shapes
sp_tree = shapes._spTree

# keep only the pictures of the template (three blue rules, two logos)
for shp in list(shapes):
    if shp.shape_type != 13:
        shp._element.getparent().remove(shp._element)

# --------------------------------------------------------------- measurement
_fonts = {}


def _font(size_pt):
    key = int(size_pt * 4)
    if key not in _fonts:
        _fonts[key] = ImageFont.truetype(ARIAL, key)
    return _fonts[key]


def n_lines(s, width_in, size_pt):
    f = _font(size_pt)
    maxw = width_in * 72 * 4 / SAFETY
    lines, cur = 1, ""
    for w in s.split():
        test = w if not cur else cur + " " + w
        if f.getlength(test) <= maxw:
            cur = test
        else:
            lines, cur = lines + 1, w
    return lines


def plain(s):
    return s.replace("**", "").replace("`", "")


# ---------------------------------------------------- the PPO loss equations
def _make_loss_figure(path, width_in=11.0):
    """Render the four PPO loss terms to a transparent PNG (matplotlib)."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    plt.rcParams["mathtext.fontset"] = "cm"

    rows = [
        (r"$\rho_t=\pi_\theta(a_t\mid s_t)\,/\,"
         r"\pi_{\theta_{\mathrm{old}}}(a_t\mid s_t)$",
         "probability ratio", 0.60),
        (r"$\mathcal{L}^{\pi}=-\,\mathbb{E}_t\!\left[\min\!\left(\rho_t\hat{A}_t,"
         r"\ \mathrm{clip}(\rho_t,1-\epsilon,1+\epsilon)\,\hat{A}_t\right)\right]$",
         "clipped policy loss", 0.66),
        (r"$\mathcal{L}^{V}=\mathbb{E}_t\!\left[(V_\theta(s_t)-\hat{R}_t)^2\right]$",
         "value loss", 0.60),
        (r"$\mathcal{L}^{H}=\mathbb{E}_t\!\left[\mathcal{H}"
         r"(\pi_\theta(\cdot\mid s_t))\right]$",
         "entropy bonus", 0.60),
        (r"$\mathcal{L}=\mathcal{L}^{\pi}+c_v\,\mathcal{L}^{V}"
         r"-c_e\,\mathcal{L}^{H}$",
         "what is backpropagated", 0.52),
    ]
    h = sum(r[2] for r in rows) + 0.10
    fig = plt.figure(figsize=(width_in, h))
    y = h - 0.06
    for eq, lab, dy in rows:
        fig.text(0.008, y / h, eq, fontsize=25, va="top", color="#202124")
        fig.text(0.995, (y - 0.10) / h, lab, fontsize=15, va="top", ha="right",
                 color="#5F6368")
        y -= dy
    fig.savefig(path, dpi=200, transparent=True)
    plt.close(fig)


# ------------------------------------------------------------------- drawing
def box(x, y, w, h, fill=WHITE, line=None, lw=1.5, shape=MSO_SHAPE.RECTANGLE):
    s = shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return s


def _emit(p, s, size, color, bold, font):
    """Emit one paragraph fragment, honouring `mono` spans inside it."""
    for piece in re.split(r"(`[^`]+`)", s):
        if not piece:
            continue
        f = font
        if len(piece) > 1 and piece.startswith("`") and piece.endswith("`"):
            piece, f = piece[1:-1], MONO
        r = p.add_run()
        r.text = piece
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = f
        r.font.color.rgb = color


def _runs(p, s, size, color, bold, font, space):
    """Mini-markup: **bold** (may contain `mono`), `mono`."""
    p.line_spacing = space   # LH (1.30) is IBM Plex Sans's natural line height
    for piece in re.split(r"(\*\*.+?\*\*)", s):
        if not piece:
            continue
        if piece.startswith("**") and piece.endswith("**"):
            _emit(p, piece[2:-2], size, color, True, font)
        else:
            _emit(p, piece, size, color, bold, font)


def text(x, y, w, lines, size=30, color=DARK, space=1.0, align=PP_ALIGN.LEFT,
         font=SANS, bold=False, indent=0.0):
    """lines: list of str or (str, overrides). Returns the height consumed."""
    total = 0.0
    for item in lines:
        o = item[1] if isinstance(item, tuple) else {}
        s = item[0] if isinstance(item, tuple) else item
        sz = o.get("size", size)
        total += o.get("gap", 0.0)
        total += n_lines(plain(s), w - (indent if s.startswith("• ") else 0),
                         sz) * sz * LH * o.get("space", space) / 72
    tb = shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(total + 0.1))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for item in lines:
        o = item[1] if isinstance(item, tuple) else {}
        s = item[0] if isinstance(item, tuple) else item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = o.get("align", align)
        if o.get("gap"):
            p.space_before = Pt(o["gap"] * 72)
        _runs(p, s, o.get("size", size), o.get("color", color),
              o.get("bold", bold), o.get("font", font), o.get("space", space))
        if s.startswith("• ") and indent:
            pPr = p._p.get_or_add_pPr()
            pPr.set("indent", str(-int(indent * 914400)))
            pPr.set("marL", str(int(indent * 914400)))
    return total


def rule(x, y, w, color=BLUE, lw=2.0):
    r = box(x, y, w, Emu(int(Pt(lw))) / 914400, color, None)
    return r


class Section:
    """A panel whose height is derived from the content put into it."""

    PAD_X, PAD_TOP, PAD_BOT, HEAD = 0.42, 0.28, 0.34, 1.08

    def __init__(self, x, w, num, title):
        self.x, self.w, self.num, self.title = x, w, num, title
        self.cx = x + self.PAD_X
        self.cw = w - 2 * self.PAD_X
        self.y = 0.0                       # provisional; shifted later
        self.cur = self.HEAD + self.PAD_TOP
        self.start = len(sp_tree)          # first content element index

    # -- content helpers, all relative to the panel top ----------------------
    def text(self, lines, **kw):
        h = text(self.cx, self.cur, kw.pop("width", None) or self.cw, lines, **kw)
        self.cur += h
        return h

    def head(self, s):
        self.cur += 0.10
        self.text([s], size=34, color=BLUE, bold=True)
        self.cur += 0.16

    def gap(self, g):
        self.cur += g

    def pic(self, path, width=None, center=True):
        width = width or self.cw
        x = self.cx + (self.cw - width) / 2 if center else self.cx
        p = shapes.add_picture(path, Inches(x), Inches(self.cur), width=Inches(width))
        self.cur += p.height / 914400
        return p

    def caption(self, s, size=23, width=None):
        self.text([s], size=size, color=MUTED, space=1.05, width=width)

    def tint(self, lines, size=27, color=BLUE, pad=0.30):
        """A tinted callout box sized to its text."""
        h = 0.0
        for item in lines:
            o = item[1] if isinstance(item, tuple) else {}
            s = item[0] if isinstance(item, tuple) else item
            sz = o.get("size", size)
            h += o.get("gap", 0.0) + n_lines(plain(s), self.cw - 2 * pad, sz) \
                * sz * LH * 1.05 / 72
        box(self.cx, self.cur, self.cw, h + 2 * pad, TINT, TINT)
        text(self.cx + pad, self.cur + pad, self.cw - 2 * pad, lines,
             size=size, color=color, space=1.05)
        self.cur += h + 2 * pad

    # -- finish --------------------------------------------------------------
    def close(self):
        self.h = self.cur + self.PAD_BOT
        frame = box(self.x, 0, self.w, self.h, WHITE, BORDER, 1.25)
        bar = box(self.x, 0, self.w, self.HEAD, BLUE, BLUE)
        tf = bar.text_frame
        tf.margin_left = Inches(1.42)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _runs(p, self.title, 46, WHITE, True, SANS, 1.0)
        badge = box(self.x + 0.17, 0.115, 0.85, 0.85, GREY, GREY)
        bt = badge.text_frame
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = bt.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        _runs(bp, str(self.num), 42, BLUE, True, SANS, 1.0)
        # frame + header behind the content
        for el in (badge._element, bar._element, frame._element):
            el.getparent().remove(el)
            sp_tree.insert(self.start, el)
        self.els = list(sp_tree)[self.start:]
        return self

    def place(self, y):
        for el in self.els:
            xfrm = el.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}off")
            if xfrm is not None:
                xfrm.set("y", str(int(xfrm.get("y")) + int(y * 914400)))
        self.y = y
        return y + self.h


def _svg_to_png(name, scale=4.0):
    """Rasterise figures/<name>.svg to figures/<name>.png (PowerPoint cannot
    place an SVG through python-pptx). Skipped if the PNG is already newer."""
    svg, png = FIG + name + ".svg", FIG + name + ".png"
    if os.path.exists(png) and os.path.getmtime(png) >= os.path.getmtime(svg):
        return png
    import cairosvg
    cairosvg.svg2png(url=svg, write_to=png, scale=scale, background_color="white")
    return png


_svg_to_png("model_architecture_poster")
_make_loss_figure(FIG + "ppo_losses.png", width_in=11.0)

# ================================================================== header ===
text(0.30, 1.85, 22.4, ["Memory Reach: GRU vs MLP under PPO on MiniGrid Memory"],
     size=62, color=DARK, bold=True, space=0.96)
text(0.32, 4.15, 22.4, ["Anh Khoa Pham"], size=44, bold=True)

# QR to the repository, in the free strip under the logos
shapes.add_picture(FIG + "qr_repo.png", Inches(30.55), Inches(3.42),
                   width=Inches(2.30))
text(22.90, 3.98, 7.45,
     ["Code, studies and", "winning checkpoints"],
     size=27, color=MUTED, align=PP_ALIGN.RIGHT, space=1.06)

LX, LW = 0.28, 12.40
RX, RW = 12.99, 19.77
TOP, BOT = 6.92, 46.52

# ========================================================== 1  Motivation ====
s1 = Section(LX, LW, 1, "Motivation")
s1.text(["• **The task is a POMDP.** The agent never sees the state, only a "
         "small egocentric patch, so the optimal action depends on information "
         "that has already left the observation.",
         ("• **The remedy is a recurrent policy** carrying a hidden state `h` "
          "forward, so an observation from a hundred steps ago can still be "
          "present at the decision.", {"gap": 0.14}),
         ("• **The concession is truncation:** full BPTT costs memory linear in "
          "sequence length, so the backward pass is cut to `L` steps — a value "
          "usually inherited in silence from whatever codebase one starts from.",
          {"gap": 0.14}),
         ("• **A memoryless policy is capped at chance (0.5)** here, so the gap "
          "from 0.5 to 1.0 success is exactly the memory a run has acquired.",
          {"gap": 0.14}),
         ("• **One PPO pipeline stays fixed;** only the encoder changes (MLP or "
          "GRU) and, for the GRU, the reach `L` ∈ {1, 4, 8, full BPTT}.",
          {"gap": 0.14})],
        size=23, space=1.06, indent=0.30)
s1.gap(0.20)
s1.tint([("**Core question:** how far back does a policy's gradient need to "
          "reach before it can solve a task a memoryless policy can't — and "
          "what does that reach cost in training time and memory?")],
        size=26, pad=0.26)
s1.gap(0.16)
s1.text(["5 studies × 50 trials × 5 seeds = **250 training runs**, ≈ 236 "
         "GPU-hours on one RTX 3060 Laptop (12 GB)."],
        size=22, color=MUTED, space=1.05)
s1.close()

# ===================================================== 2  Problem setting ====
s2 = Section(RX, RW, 2, "Problem Setting")
s2.head("The environment")
s2.text(["• `MiniGrid-MemoryS17Random-v0`: a **cue** (a key or a ball) sits in a "
         "room at the west end; the corridor forks into a **T-junction** at the "
         "east end with a key at one prong and a ball at the other. The episode "
         "succeeds **only** if the agent walks to the prong matching the cue. "
         "Reward is `1 − 0.9 · steps / max_steps` on success, 0 otherwise.",
         ("• **The agent only sees a partial view in front of it:** a 7 × 7 "
          "egocentric patch of what lies ahead (7 × 7 × 3 → one-hot 980 "
          "features, 7 discrete actions). Walls block the view.", {"gap": 0.13}),
         ("• **The cue is out of sight when the decision is made.** At the "
          "junction the observation is identical whichever answer is correct, "
          "so a memoryless policy cannot beat chance however well it is trained.",
          {"gap": 0.13}),
         ("• **The information has to be sought out.** The agent spawns at a "
          "random corridor position facing east, away from the cue, so it must "
          "turn around, walk back, look and turn again — a detour that costs "
          "reward at once and pays nothing until the memory also works.",
          {"gap": 0.13})],
        size=23, space=1.06, indent=0.30)
s2.gap(0.26)
s2.head("The two encoders")
s2.text(["• **MLP — memoryless (the baseline).** 1–4 Linear+ReLU blocks applied "
         "to each timestep independently: no state crosses the sequence and its "
         "gradient spans a single step. Whatever it scores is the ceiling for a "
         "policy that can only react to the frame in front of it.",
         ("• **GRU — with memory.** The same input plus a single recurrent layer "
          "of tuned width carrying `h` forward; the gradient flows backwards "
          "along that chain and `L` caps how many steps back it may reach. "
          "Everything else — PPO, rollout, budget, evaluation — is identical.",
          {"gap": 0.14})],
        size=23, space=1.06, indent=0.30)
s2.gap(0.28)
s2.pic(FIG + "task_memory_s17.png", width=8.6)
s2.gap(0.10)
s2.caption("The cue sits in the west room; the agent spawns in the corridor "
           "facing away from it, so the information has to be sought out.",
           size=22)
s2.close()

# ============================================================ 3  Approach ====
s3 = Section(LX, LW, 3, "Approach")
s3.head("Network")
s3.pic(FIG + "model_architecture_poster.png", width=10.0)
s3.gap(0.06)
s3.caption("7 × 7 × 3 view → one-hot (980) → encoder → linear actor + linear "
           "critic. The encoder is the **only** difference between arms.",
           size=22)
s3.gap(0.12)
s3.head("PPO objective")
s3.text(["• One iteration collects W = 32 workers × T = 361 steps and computes "
         "**GAE advantages over the whole rollout**, before any split, with "
         "returns R̂ = Â + V(s).",
         ("• The rollout is then cut at every episode boundary **and every `L` "
          "steps**, each chunk seeded with the hidden state recorded there, so "
          "**only the backward pass is shortened**.", {"gap": 0.13}),
         ("• Three shuffled epochs minimise a clipped policy loss, a value loss "
          "and an entropy bonus:", {"gap": 0.13})],
        size=23, space=1.06, indent=0.30)
s3.gap(0.08)
s3.pic(FIG + "ppo_losses.png", width=9.8)
s3.gap(0.06)
s3.caption("Every expectation covers real timesteps only: chunks are padded "
           "inside a minibatch and padding never enters the loss; the gradient "
           "norm is clipped before each Adam step.", size=22)
s3.gap(0.12)
s3.head("Tuning protocol")
s3.text(["• One Optuna study **per arm**: 50 TPE trials over **identical** "
         "ranges — no arm is judged on another's settings.",
         ("• Each trial trains the same 5 seeds for 1000 iterations and is "
          "scored **mean − std across seeds** on 50 held-out mazes.",
          {"gap": 0.12}),
         ("• Nothing that buys compute is searchable; the minibatch size is "
          "set by a VRAM probe, so every arm is measured at the **same "
          "memory budget**.", {"gap": 0.12})],
        size=23, space=1.06, indent=0.28)
s3.close()

# ============================================================= 4  Results ====
s4 = Section(RX, RW, 4, "Results")
cols = [4.45, 3.05, 3.05, 2.05, 2.45, 1.88]
hdr = ["arm", "return", "success rate", "steps", "minibatch", "time"]
rows = [
    ["MLP (memoryless)", "0.556 ± 0.029", "0.592 ± 0.027", "24.3", "1024", "1.3 h"],
    ["GRU, L = 1", "0.573 ± 0.064", "0.588 ± 0.068", "9.8", "4096", "2.6 h"],
    ["GRU, L = 4", "0.492 ± 0.001", "0.500 ± 0.000", "7.0", "4096", "1.6 h"],
    ["GRU, L = 8", "0.958 ± 0.004", "1.000 ± 0.000", "16.8", "4096", "2.0 h"],
    ["GRU, full BPTT", "0.958 ± 0.002", "1.000 ± 0.000", "16.8", "512", "1.3 h"],
]
WIN, RH = {3, 4}, 0.60
rule(s4.cx, s4.cur, s4.cw, BLUE, 3)
xo = s4.cx
for j, c in enumerate(hdr):
    text(xo + 0.06, s4.cur + 0.17, cols[j] - 0.12, [c], size=25, color=BLUE,
         bold=True, align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    xo += cols[j]
s4.cur += 0.82
rule(s4.cx, s4.cur, s4.cw, BLUE, 1.5)
s4.cur += 0.10
for i, r in enumerate(rows):
    if i in WIN:
        box(s4.cx, s4.cur - 0.04, s4.cw, RH, TINT, TINT)
    xo = s4.cx
    for j, c in enumerate(r):
        hit = (i in WIN and j in (1, 2)) or (i == 4 and j == 5)
        text(xo + 0.06, s4.cur + 0.06, cols[j] - 0.12, [c], size=25,
             color=GREEN if hit else DARK, bold=hit,
             align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER,
             font=SANS if j == 0 else MONO)
        xo += cols[j]
    s4.cur += RH
s4.cur += 0.04
rule(s4.cx, s4.cur, s4.cw, BLUE, 3)
s4.cur += 0.16
s4.caption("Winning trial per study, averaged over its five seeds; ± is the "
           "spread across seeds. No setting produced an intermediate score.",
           size=22)
s4.gap(0.38)
s4.pic(FIG + "compare_curves.png", width=12.3)
s4.gap(0.12)
s4.caption("Mean over 5 seeds, band at ± 1 std, dashed line at chance. The "
           "two solving arms separate before iteration 400; the three failing "
           "arms never leave the chance line.", size=22)
s4.close()

# ======================================================== 5  Key insights ====
s5 = Section(LX, LW, 5, "Key Insights")
s5.text(["• **Carrying `h` forward is not what matters.** Every GRU arm "
         "does that, including the ones that fail. What separates them is "
         "whether the **gradient** reaches back far enough; past 8 steps "
         "adds nothing.",
         ("• **What must fit inside `L` is the gap, not the episode.** The "
          "gap runs from the last sighting of the cue to the choice at the "
          "fork: ≤ 8 steps in **62 %** of mazes by BFS, ≤ 4 in only "
          "**26 %**.", {"gap": 0.06}),
         ("• **The failing arms are not partially remembering.** At 7.0 "
          "steps the L = 4 arm rushes to the fork and always picks the "
          "same side, 25 of 50 here: its std of 0.000 is **arithmetic, "
          "not convergence**.", {"gap": 0.06}),
         ("• **The reach costs memory, not time.** Full BPTT is the "
          "**fastest** GRU arm (1.3 h vs 2.6 h at L = 1); halving `L` "
          "doubles the padded sequences. What it pays instead is VRAM: "
          "512 vs 4096.", {"gap": 0.06})],
        size=23, space=1.06, indent=0.30)
s5.gap(0.08)
s5.tint(["**Recommendation.** On this task: full BPTT, with L = 8 as the "
         "fallback on a smaller GPU. Below 8 you do not get a cheaper "
         "solution; you get the MLP's guessing policy at GRU prices."],
        size=24, pad=0.12)
s5.close()

# =========================================================== 6  Future work ==
s6 = Section(RX, RW, 6, "Limitations & Future Work")
s6.text(["• **One task, one architecture, one algorithm.** `L` ∈ {5,6,7} was "
         "never run, so the threshold is **bracketed, not located**, and five "
         "seeds separate 0.5 from 1.0 safely but support no finer claim.",
         ("• **Locate the threshold.** Sweep `L` ∈ {5,6,7} and set the "
          "corridor length directly instead of leaving it to the "
          "environment: a causal test of the gap account.", {"gap": 0.28}),
         ("• **Add R2D2-style burn-in.** Does a short window then behave "
          "like a long one? If so the threshold is about optimisation, not "
          "the gradient path.", {"gap": 0.22}),
         ("• **Repeat with an LSTM and a small transformer**, where the "
          "analogous knob is the attention window.", {"gap": 0.20})],
        size=25, space=1.08, indent=0.30)
s6.gap(0.40)
rule(s6.cx, s6.cur, s6.cw, RGBColor(0xD5, 0xD8, 0xDC), 1.5)
s6.gap(0.24)
s6.text(["**References.** Hausknecht & Stone 2015 · Kapturowski et al. 2019 "
         "(R2D2) · Ni et al. 2022 · Schulman et al. 2017 (PPO) · "
         "Chevalier-Boisvert et al. 2023 (MiniGrid)"],
        size=22, color=MUTED, space=1.08)
s6.close()

# ============================================================ distribute =====
for col, secs in (("left", [s1, s3, s5]), ("right", [s2, s4, s6])):
    used = sum(s.h for s in secs)
    slack = (BOT - TOP) - used
    gap = slack / (len(secs) - 1)
    print(f"{col:5s} used={used:6.2f} slack={slack:6.2f} gap={gap:5.2f} "
          f"[{' '.join(f'{s.num}:{s.h:.2f}' for s in secs)}]")
    if slack < 0:
        print(f"  !! {col} column overflows by {-slack:.2f} in")
    y = TOP
    for s in secs:
        y = s.place(y) + gap

# (the template's rules and logos are already at the bottom of the z-order:
# they are the only shapes that survived the clean-up, and everything drawn
# afterwards was appended after them.)

prs.save(OUT)
print("written", OUT)
